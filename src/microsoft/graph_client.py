"""Microsoft Graph API wrapper."""

import io
import logging
from datetime import datetime, timedelta, timezone
from typing import Any

import httpx

logger = logging.getLogger(__name__)


def _extract_docx_text(content: bytes) -> str:
    """Extract text from a .docx file."""
    from docx import Document
    doc = Document(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_xlsx_text(content: bytes) -> str:
    """Extract text from a .xlsx file."""
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    result = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        result.append(f"=== Sheet: {sheet_name} ===")
        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell) if cell is not None else "" for cell in row]
            if any(v.strip() for v in row_values):
                result.append("\t".join(row_values))
    wb.close()
    return "\n".join(result)


def _extract_pptx_text(content: bytes) -> str:
    """Extract text from a .pptx file."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    result = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        if slide_text:
            result.append(f"=== Slide {i} ===")
            result.extend(slide_text)
    return "\n\n".join(result)


def _extract_pdf_text(content: bytes) -> str:
    """Extract text from a PDF file."""
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(content))
    result = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            result.append(f"=== Page {i} ===")
            result.append(text)
    return "\n\n".join(result)

GRAPH_BASE_URL = "https://graph.microsoft.com/v1.0"


class GraphClient:
    """Client for Microsoft Graph API."""

    def __init__(self, access_token: str) -> None:
        self.access_token = access_token
        self.headers = {
            "Authorization": f"Bearer {access_token}",
            "Content-Type": "application/json",
        }

    async def _request(
        self,
        method: str,
        endpoint: str,
        params: dict | None = None,
        json_data: dict | None = None,
    ) -> dict[str, Any]:
        """Make a request to the Graph API."""
        url = f"{GRAPH_BASE_URL}{endpoint}"

        async with httpx.AsyncClient() as client:
            response = await client.request(
                method=method,
                url=url,
                headers=self.headers,
                params=params,
                json=json_data,
                timeout=30.0,
            )

            if response.status_code == 401:
                raise PermissionError("Access token expired or invalid")
            elif response.status_code == 403:
                raise PermissionError("Insufficient permissions for this operation")
            elif response.status_code >= 400:
                error_data = response.json() if response.content else {}
                error_msg = error_data.get("error", {}).get("message", response.text)
                raise Exception(f"Graph API error ({response.status_code}): {error_msg}")

            return response.json() if response.content else {}

    # ==================== EMAIL ====================

    async def get_emails(
        self,
        limit: int = 10,
        search: str | None = None,
        folder: str = "inbox",
    ) -> list[dict]:
        """Get recent emails, optionally filtered by search query."""
        params = {
            "$top": limit,
            "$select": "id,subject,from,receivedDateTime,bodyPreview,isRead,importance",
        }

        if search:
            # Note: $orderby is not supported with $search
            params["$search"] = f'"{search}"'
        else:
            params["$orderby"] = "receivedDateTime desc"

        endpoint = f"/me/mailFolders/{folder}/messages"
        result = await self._request("GET", endpoint, params=params)

        emails = []
        for msg in result.get("value", []):
            emails.append({
                "id": msg["id"],
                "subject": msg.get("subject", "(No subject)"),
                "from": msg.get("from", {}).get("emailAddress", {}).get("address", "Unknown"),
                "from_name": msg.get("from", {}).get("emailAddress", {}).get("name", ""),
                "received": msg.get("receivedDateTime", ""),
                "preview": msg.get("bodyPreview", "")[:200],
                "is_read": msg.get("isRead", False),
                "importance": msg.get("importance", "normal"),
            })

        return emails

    async def get_email(self, email_id: str) -> dict:
        """Get a specific email by ID."""
        params = {
            "$select": "id,subject,from,toRecipients,receivedDateTime,body,isRead,importance,hasAttachments",
        }

        result = await self._request("GET", f"/me/messages/{email_id}", params=params)

        return {
            "id": result["id"],
            "subject": result.get("subject", "(No subject)"),
            "from": result.get("from", {}).get("emailAddress", {}).get("address", "Unknown"),
            "from_name": result.get("from", {}).get("emailAddress", {}).get("name", ""),
            "to": [r.get("emailAddress", {}).get("address", "") for r in result.get("toRecipients", [])],
            "received": result.get("receivedDateTime", ""),
            "body": result.get("body", {}).get("content", ""),
            "body_type": result.get("body", {}).get("contentType", "text"),
            "is_read": result.get("isRead", False),
            "importance": result.get("importance", "normal"),
            "has_attachments": result.get("hasAttachments", False),
        }

    # ==================== CALENDAR ====================

    async def get_calendar_events(
        self,
        days: int = 7,
        past_days: int = 0,
        limit: int = 50,
    ) -> list[dict]:
        """
        Get calendar events within a date range.

        Args:
            days: Number of days to look ahead (default 7)
            past_days: Number of days to look back (default 0)
            limit: Maximum number of events to return
        """
        now = datetime.now(timezone.utc)
        start_dt = now - timedelta(days=past_days)
        end_dt = now + timedelta(days=days)

        params = {
            "startDateTime": start_dt.isoformat(),
            "endDateTime": end_dt.isoformat(),
            "$orderby": "start/dateTime",
            "$select": "id,subject,start,end,location,organizer,attendees,isOnlineMeeting,onlineMeetingUrl,bodyPreview",
            "$top": limit,
        }

        result = await self._request("GET", "/me/calendarView", params=params)

        events = []
        for event in result.get("value", []):
            start = event.get("start", {})
            end_time = event.get("end", {})

            events.append({
                "id": event["id"],
                "subject": event.get("subject", "(No title)"),
                "start": start.get("dateTime", ""),
                "start_timezone": start.get("timeZone", "UTC"),
                "end": end_time.get("dateTime", ""),
                "end_timezone": end_time.get("timeZone", "UTC"),
                "location": event.get("location", {}).get("displayName", ""),
                "organizer": event.get("organizer", {}).get("emailAddress", {}).get("name", ""),
                "organizer_email": event.get("organizer", {}).get("emailAddress", {}).get("address", ""),
                "attendees": [
                    {
                        "name": a.get("emailAddress", {}).get("name", ""),
                        "email": a.get("emailAddress", {}).get("address", ""),
                        "status": a.get("status", {}).get("response", ""),
                    }
                    for a in event.get("attendees", [])
                ],
                "is_online": event.get("isOnlineMeeting", False),
                "online_url": event.get("onlineMeetingUrl", ""),
                "description": event.get("bodyPreview", "")[:200],
            })

        return events

    async def get_past_events(self, days: int = 7, limit: int = 50) -> list[dict]:
        """Get past calendar events from the last N days."""
        return await self.get_calendar_events(days=0, past_days=days, limit=limit)

    async def get_events_for_date(self, date_str: str, limit: int = 50) -> list[dict]:
        """
        Get calendar events for a specific date.

        Args:
            date_str: Date in YYYY-MM-DD format (e.g., "2025-01-30")
            limit: Maximum number of events
        """
        from datetime import datetime as dt

        try:
            target_date = dt.strptime(date_str, "%Y-%m-%d").replace(tzinfo=timezone.utc)
        except ValueError:
            raise ValueError(f"Invalid date format: {date_str}. Use YYYY-MM-DD format.")

        start_of_day = target_date.replace(hour=0, minute=0, second=0, microsecond=0)
        end_of_day = target_date.replace(hour=23, minute=59, second=59, microsecond=999999)

        params = {
            "startDateTime": start_of_day.isoformat(),
            "endDateTime": end_of_day.isoformat(),
            "$orderby": "start/dateTime",
            "$select": "id,subject,start,end,location,organizer,attendees,isOnlineMeeting,onlineMeetingUrl,bodyPreview",
            "$top": limit,
        }

        result = await self._request("GET", "/me/calendarView", params=params)

        events = []
        for event in result.get("value", []):
            start = event.get("start", {})
            end_time = event.get("end", {})

            events.append({
                "id": event["id"],
                "subject": event.get("subject", "(No title)"),
                "start": start.get("dateTime", ""),
                "start_timezone": start.get("timeZone", "UTC"),
                "end": end_time.get("dateTime", ""),
                "end_timezone": end_time.get("timeZone", "UTC"),
                "location": event.get("location", {}).get("displayName", ""),
                "organizer": event.get("organizer", {}).get("emailAddress", {}).get("name", ""),
                "organizer_email": event.get("organizer", {}).get("emailAddress", {}).get("address", ""),
                "attendees": [
                    {
                        "name": a.get("emailAddress", {}).get("name", ""),
                        "email": a.get("emailAddress", {}).get("address", ""),
                        "status": a.get("status", {}).get("response", ""),
                    }
                    for a in event.get("attendees", [])
                ],
                "is_online": event.get("isOnlineMeeting", False),
                "online_url": event.get("onlineMeetingUrl", ""),
                "description": event.get("bodyPreview", "")[:200],
            })

        return events

    async def get_today_events(self) -> list[dict]:
        """Get today's calendar events."""
        now = datetime.now(timezone.utc)
        start_of_day = now.replace(hour=0, minute=0, second=0, microsecond=0)
        end_of_day = start_of_day + timedelta(days=1)

        params = {
            "startDateTime": start_of_day.isoformat(),
            "endDateTime": end_of_day.isoformat(),
            "$orderby": "start/dateTime",
            "$select": "id,subject,start,end,location,isOnlineMeeting,onlineMeetingUrl",
        }

        result = await self._request("GET", "/me/calendarView", params=params)

        events = []
        for event in result.get("value", []):
            start = event.get("start", {})
            end = event.get("end", {})

            events.append({
                "id": event["id"],
                "subject": event.get("subject", "(No title)"),
                "start": start.get("dateTime", ""),
                "end": end.get("dateTime", ""),
                "location": event.get("location", {}).get("displayName", ""),
                "is_online": event.get("isOnlineMeeting", False),
                "online_url": event.get("onlineMeetingUrl", ""),
            })

        return events

    # ==================== TEAMS ====================

    async def get_teams_chats(self, limit: int = 10) -> list[dict]:
        """Get recent Teams chats."""
        params = {
            "$top": limit,
            "$orderby": "lastMessagePreview/createdDateTime desc",
            "$expand": "lastMessagePreview",
        }

        result = await self._request("GET", "/me/chats", params=params)

        chats = []
        for chat in result.get("value", []):
            last_msg = chat.get("lastMessagePreview", {})

            chats.append({
                "id": chat["id"],
                "topic": chat.get("topic", ""),
                "chat_type": chat.get("chatType", ""),
                "last_message": last_msg.get("body", {}).get("content", "")[:100] if last_msg else "",
                "last_message_from": last_msg.get("from", {}).get("user", {}).get("displayName", "") if last_msg else "",
                "last_message_time": last_msg.get("createdDateTime", "") if last_msg else "",
            })

        return chats

    async def get_chat_messages(self, chat_id: str, limit: int = 20) -> list[dict]:
        """Get messages from a Teams chat."""
        params = {
            "$top": limit,
            "$orderby": "createdDateTime desc",
        }

        result = await self._request("GET", f"/me/chats/{chat_id}/messages", params=params)

        messages = []
        for msg in result.get("value", []):
            from_user = msg.get("from", {})
            user_info = from_user.get("user", {}) if from_user else {}

            messages.append({
                "id": msg["id"],
                "content": msg.get("body", {}).get("content", ""),
                "content_type": msg.get("body", {}).get("contentType", "text"),
                "from": user_info.get("displayName", "Unknown"),
                "from_email": user_info.get("email", ""),
                "created": msg.get("createdDateTime", ""),
                "message_type": msg.get("messageType", ""),
            })

        return messages

    # ==================== FILES ====================

    async def search_files(self, query: str, limit: int = 10) -> list[dict]:
        """Search files in OneDrive and SharePoint."""
        # Use the search API
        search_body = {
            "requests": [
                {
                    "entityTypes": ["driveItem"],
                    "query": {"queryString": query},
                    "from": 0,
                    "size": limit,
                }
            ]
        }

        result = await self._request("POST", "/search/query", json_data=search_body)

        files = []
        for response in result.get("value", []):
            for hit in response.get("hitsContainers", [{}])[0].get("hits", []):
                resource = hit.get("resource", {})
                parent_ref = resource.get("parentReference", {})
                files.append({
                    "id": resource.get("id", ""),
                    "drive_id": parent_ref.get("driveId", ""),
                    "name": resource.get("name", ""),
                    "web_url": resource.get("webUrl", ""),
                    "size": resource.get("size", 0),
                    "created": resource.get("createdDateTime", ""),
                    "modified": resource.get("lastModifiedDateTime", ""),
                    "created_by": resource.get("createdBy", {}).get("user", {}).get("displayName", ""),
                })

        return files

    async def get_recent_files(self, limit: int = 10) -> list[dict]:
        """Get recently accessed files."""
        params = {
            "$top": limit,
            "$orderby": "lastAccessedDateTime desc",
        }

        result = await self._request("GET", "/me/drive/recent", params=params)

        files = []
        for item in result.get("value", []):
            files.append({
                "id": item.get("id", ""),
                "name": item.get("name", ""),
                "web_url": item.get("webUrl", ""),
                "size": item.get("size", 0),
                "modified": item.get("lastModifiedDateTime", ""),
            })

        return files

    async def get_file_content(self, file_id: str, drive_id: str | None = None, max_size_mb: float = 5.0) -> dict:
        """
        Download and return file content.

        Args:
            file_id: The ID of the file to download
            drive_id: The drive ID (required for SharePoint files, optional for OneDrive)
            max_size_mb: Maximum file size to download in MB

        Returns:
            Dict with file metadata and content
        """
        # Build the correct endpoint based on whether we have a drive_id
        if drive_id:
            # SharePoint or specific drive
            base_path = f"/drives/{drive_id}/items/{file_id}"
        else:
            # Default to user's OneDrive
            base_path = f"/me/drive/items/{file_id}"

        logger.info(f"get_file_content: base_path={base_path}")

        # First get file metadata
        metadata = await self._request("GET", base_path)

        file_name = metadata.get("name", "")
        file_size = metadata.get("size", 0)
        mime_type = metadata.get("file", {}).get("mimeType", "")
        logger.info(f"get_file_content: file={file_name}, size={file_size}, mime={mime_type}")

        # Check file size
        max_size_bytes = int(max_size_mb * 1024 * 1024)
        if file_size > max_size_bytes:
            return {
                "error": f"File too large ({file_size / 1024 / 1024:.1f} MB). Max size is {max_size_mb} MB.",
                "name": file_name,
                "size": file_size,
            }

        # Determine file type and how to handle it
        extension = file_name.lower().split(".")[-1] if "." in file_name else ""

        text_extensions = {"txt", "md", "csv", "json", "xml", "html", "css", "js", "ts", "py", "yaml", "yml", "log", "ini", "cfg"}
        office_extensions = {"docx", "xlsx", "pptx", "doc", "xls", "ppt"}

        content_url = f"{GRAPH_BASE_URL}{base_path}/content"

        async with httpx.AsyncClient() as client:
            if extension in text_extensions:
                # Download raw content for text files
                response = await client.get(
                    content_url,
                    headers=self.headers,
                    follow_redirects=True,
                    timeout=60.0,
                )

                if response.status_code == 200:
                    try:
                        content = response.text
                        return {
                            "name": file_name,
                            "size": file_size,
                            "mime_type": mime_type,
                            "extension": extension,
                            "content": content[:50000],  # Limit content size
                            "truncated": len(content) > 50000,
                        }
                    except Exception as e:
                        return {"error": f"Failed to decode file content: {e}", "name": file_name}
                else:
                    return {"error": f"Failed to download file: {response.status_code}", "name": file_name}

            elif extension in office_extensions or extension == "pdf":
                # Download the file and extract text using appropriate parser
                logger.info(f"get_file_content: downloading {extension} file from {content_url}")
                response = await client.get(
                    content_url,
                    headers=self.headers,
                    follow_redirects=True,
                    timeout=60.0,
                )

                logger.info(f"get_file_content: download response status={response.status_code}")
                if response.status_code != 200:
                    logger.error(f"get_file_content: download failed with {response.status_code}: {response.text[:500]}")
                    return {"error": f"Failed to download file: {response.status_code}", "name": file_name}

                file_bytes = response.content
                logger.info(f"get_file_content: downloaded {len(file_bytes)} bytes, extracting text...")

                try:
                    if extension == "docx":
                        content = _extract_docx_text(file_bytes)
                        logger.info(f"get_file_content: extracted {len(content)} chars from docx")
                    elif extension == "xlsx":
                        content = _extract_xlsx_text(file_bytes)
                        logger.info(f"get_file_content: extracted {len(content)} chars from xlsx")
                    elif extension == "pptx":
                        content = _extract_pptx_text(file_bytes)
                        logger.info(f"get_file_content: extracted {len(content)} chars from pptx")
                    elif extension == "pdf":
                        content = _extract_pdf_text(file_bytes)
                        logger.info(f"get_file_content: extracted {len(content)} chars from pdf")
                    elif extension in {"doc", "xls", "ppt"}:
                        return {
                            "name": file_name,
                            "size": file_size,
                            "mime_type": mime_type,
                            "extension": extension,
                            "web_url": metadata.get("webUrl", ""),
                            "note": f"Legacy Office format (.{extension}) not supported. Please convert to .{extension}x format.",
                        }
                    else:
                        content = ""

                    if content:
                        return {
                            "name": file_name,
                            "size": file_size,
                            "mime_type": mime_type,
                            "extension": extension,
                            "content": content[:50000],
                            "truncated": len(content) > 50000,
                        }
                    else:
                        return {
                            "name": file_name,
                            "size": file_size,
                            "mime_type": mime_type,
                            "extension": extension,
                            "web_url": metadata.get("webUrl", ""),
                            "note": "No text content could be extracted from this file.",
                        }

                except Exception as e:
                    logger.error(f"Failed to extract content from {extension} file: {e}")
                    return {
                        "error": f"Failed to extract text from {extension} file: {str(e)}",
                        "name": file_name,
                        "web_url": metadata.get("webUrl", ""),
                    }

            else:
                return {
                    "name": file_name,
                    "size": file_size,
                    "mime_type": mime_type,
                    "extension": extension,
                    "web_url": metadata.get("webUrl", ""),
                    "note": f"Unsupported file type: {extension}. Use web URL to access.",
                }

    # ==================== PERSON SEARCH ====================

    async def get_emails_from_person(
        self,
        person: str,
        limit: int = 10,
        unread_only: bool = False,
    ) -> list[dict]:
        """Get emails from a specific person by name or email address."""
        # Note: $orderby is not supported with $search, but search returns relevance-ranked results
        params = {
            "$top": limit,
            "$select": "id,subject,from,receivedDateTime,bodyPreview,isRead,importance",
            "$search": f'"from:{person}"',
        }

        if unread_only:
            params["$filter"] = "isRead eq false"

        endpoint = "/me/mailFolders/inbox/messages"
        result = await self._request("GET", endpoint, params=params)

        emails = []
        for msg in result.get("value", []):
            emails.append({
                "id": msg["id"],
                "subject": msg.get("subject", "(No subject)"),
                "from": msg.get("from", {}).get("emailAddress", {}).get("address", "Unknown"),
                "from_name": msg.get("from", {}).get("emailAddress", {}).get("name", ""),
                "received": msg.get("receivedDateTime", ""),
                "preview": msg.get("bodyPreview", "")[:200],
                "is_read": msg.get("isRead", False),
                "importance": msg.get("importance", "normal"),
            })

        return emails

    async def get_teams_messages_from_person(
        self,
        person: str,
        limit: int = 20,
        chat_type: str | None = None,
        include_context: bool = False,
    ) -> list[dict]:
        """
        Get Teams messages from a specific person across chats.

        Args:
            person: Name of the person to search for
            limit: Maximum number of messages to return
            chat_type: Filter by chat type - "oneOnOne" for 1:1, "group" for group chats, None for all
            include_context: If True, include your replies for full conversation context
        """
        # First get recent chats
        chats_params = {
            "$top": 50,
            "$orderby": "lastMessagePreview/createdDateTime desc",
        }
        chats_result = await self._request("GET", "/me/chats", params=chats_params)

        person_lower = person.lower()
        messages = []
        chats_with_person = []

        # Find chats that have messages from the person
        for chat in chats_result.get("value", []):
            current_chat_type = chat.get("chatType", "")
            # Filter by chat type if specified
            if chat_type and current_chat_type != chat_type:
                continue

            chat_id = chat["id"]
            chat_topic = chat.get("topic", "")

            try:
                msg_params = {
                    "$top": 50,
                    "$orderby": "createdDateTime desc",
                }
                msg_result = await self._request(
                    "GET", f"/me/chats/{chat_id}/messages", params=msg_params
                )

                chat_messages = msg_result.get("value", [])
                has_person_message = any(
                    person_lower in (msg.get("from", {}).get("user", {}).get("displayName", "") or "").lower()
                    for msg in chat_messages
                )

                if has_person_message:
                    chats_with_person.append({
                        "chat_id": chat_id,
                        "chat_topic": chat_topic,
                        "chat_type": current_chat_type,
                        "messages": chat_messages,
                    })

            except Exception as e:
                logger.warning(f"Failed to get messages from chat {chat_id}: {e}")
                continue

        # Process messages
        for chat_info in chats_with_person:
            for msg in chat_info["messages"]:
                from_user = msg.get("from", {})
                user_info = from_user.get("user", {}) if from_user else {}
                display_name = user_info.get("displayName", "")
                content = msg.get("body", {}).get("content", "")

                if not content or msg.get("messageType") != "message":
                    continue

                is_from_person = person_lower in display_name.lower()

                # Include message if it's from the person, or if include_context is True
                if is_from_person or include_context:
                    messages.append({
                        "id": msg["id"],
                        "chat_id": chat_info["chat_id"],
                        "chat_topic": chat_info["chat_topic"],
                        "chat_type": chat_info["chat_type"],
                        "content": content[:500],
                        "content_type": msg.get("body", {}).get("contentType", "text"),
                        "from": display_name,
                        "from_email": user_info.get("email", ""),
                        "created": msg.get("createdDateTime", ""),
                        "is_from_searched_person": is_from_person,
                    })

                if len(messages) >= limit:
                    break

            if len(messages) >= limit:
                break

        # Sort by date descending
        messages.sort(key=lambda x: x.get("created", ""), reverse=True)
        return messages[:limit]

    # ==================== USER INFO ====================

    async def get_me(self) -> dict:
        """Get the current user's profile."""
        result = await self._request("GET", "/me")

        return {
            "id": result.get("id", ""),
            "name": result.get("displayName", ""),
            "email": result.get("mail", result.get("userPrincipalName", "")),
            "job_title": result.get("jobTitle", ""),
            "department": result.get("department", ""),
            "office": result.get("officeLocation", ""),
        }
